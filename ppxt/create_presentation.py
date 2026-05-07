from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
 
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
 
# Color palette - teal/blue professional (matching Sysplus brand)
TEAL_DARK  = RGBColor(0x3A, 0x7D, 0x8C)   # #3A7D8C - header teal
TEAL_MED   = RGBColor(0x4D, 0xA6, 0xB8)   # #4DA6B8 - medium teal
TEAL_LIGHT = RGBColor(0xD6, 0xEE, 0xF3)   # #D6EEF3 - light teal bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x1E, 0x2A, 0x35)   # near-black for text
MID_GRAY   = RGBColor(0x55, 0x6B, 0x7A)   # body text gray
ORANGE_ACC = RGBColor(0xFF, 0x8C, 0x2A)   # orange accent (from logo)
LIGHT_GRAY = RGBColor(0xF4, 0xF8, 0xFA)   # slide bg
 
W = prs.slide_width
H = prs.slide_height
 
def add_rect(slide, x, y, w, h, fill_color, transparency=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape
 
def add_text_box(slide, text, x, y, w, h, font_size, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox
 
def add_slide_base(prs, bg_color=None):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    # background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color if bg_color else LIGHT_GRAY
    return slide
 
def add_header_bar(slide, title_text, subtitle=None):
    """Dark teal top bar with title"""
    bar = add_rect(slide, 0, 0, W, Inches(1.1), TEAL_DARK)
    add_text_box(slide, title_text, Inches(0.5), Inches(0.15), W - Inches(1), Inches(0.8),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
    # orange accent dot
    dot = slide.shapes.add_shape(9, Inches(0.18), Inches(0.35), Inches(0.14), Inches(0.14))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ORANGE_ACC
    dot.line.fill.background()
 
def add_footer(slide):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), TEAL_DARK)
    add_text_box(slide, "SYSPLUS  |  Marketplace Completo", 0, H - Inches(0.38), W, Inches(0.38),
                 font_size=9, color=RGBColor(0xCC, 0xE8, 0xEE), align=PP_ALIGN.CENTER, font_name="Calibri")
 
# ─────────────────────────────────────────
# SLIDE 1 — Nome do Projeto
# ─────────────────────────────────────────
slide1 = add_slide_base(prs, TEAL_DARK)
 
# Big centered background circle decoration
circ = slide1.shapes.add_shape(9, Inches(3.5), Inches(0.5), Inches(6), Inches(6))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x2E, 0x6A, 0x78)
circ.line.fill.background()
 
# Logo-style "S+" box
box = add_rect(slide1, Inches(5.4), Inches(1.6), Inches(2.4), Inches(2.4), ORANGE_ACC)
add_text_box(slide1, "S+", Inches(5.4), Inches(1.6), Inches(2.4), Inches(2.4),
             font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
 
add_text_box(slide1, "SYSPLUS", Inches(0.8), Inches(2.0), Inches(4.5), Inches(1.4),
             font_size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text_box(slide1, "Marketplace Completo", Inches(0.8), Inches(3.2), Inches(4.5), Inches(0.7),
             font_size=22, bold=False, color=TEAL_LIGHT, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text_box(slide1, "Conectando compradores e vendedores\nem uma experiência simples e confiável.",
             Inches(0.8), Inches(3.9), Inches(4.5), Inches(1.2),
             font_size=14, color=RGBColor(0xAA, 0xD4, 0xDD), align=PP_ALIGN.LEFT, font_name="Calibri")
 
# Tech tag
tag = add_rect(slide1, Inches(0.8), Inches(5.4), Inches(2.5), Inches(0.45), ORANGE_ACC)
add_text_box(slide1, "HTML  |  CSS  |  JavaScript", Inches(0.8), Inches(5.4), Inches(2.5), Inches(0.45),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
 
# ─────────────────────────────────────────
# SLIDE 2 — Integrantes
# ─────────────────────────────────────────
slide2 = add_slide_base(prs, LIGHT_GRAY)
add_header_bar(slide2, "Integrantes do Projeto")
add_footer(slide2)
 
# Section label
add_text_box(slide2, "EQUIPE DE DESENVOLVIMENTO", Inches(0.5), Inches(1.2), Inches(12), Inches(0.4),
             font_size=11, bold=True, color=TEAL_MED, align=PP_ALIGN.CENTER, font_name="Calibri")
 
members = [
    ("Artur Furtado de\nMagalhães Garcia", "Desenvolvedor"),
    ("Lucas Fernandes\nBraga Silva", "Desenvolvedor"),
    ("Guilherme Henrique\nRodrigues Faria", "Desenvolvedor"),
]
card_w = Inches(3.4)
card_h = Inches(3.6)
gap = Inches(0.45)
total_w = 3 * card_w + 2 * gap
start_x = (W - total_w) / 2
card_y = Inches(1.75)
 
for i, (name, role) in enumerate(members):
    cx = start_x + i * (card_w + gap)
    # Card shadow
    shadow = add_rect(slide2, cx + Inches(0.05), card_y + Inches(0.05), card_w, card_h, RGBColor(0xCC, 0xD8, 0xDD))
    # Card body
    card = add_rect(slide2, cx, card_y, card_w, card_h, WHITE)
    # Top accent
    accent = add_rect(slide2, cx, card_y, card_w, Inches(0.12), TEAL_DARK)
    # Avatar circle
    av = slide2.shapes.add_shape(9, cx + card_w/2 - Inches(0.65), card_y + Inches(0.3), Inches(1.3), Inches(1.3))
    av.fill.solid()
    av.fill.fore_color.rgb = TEAL_LIGHT
    av.line.solid()
    av.line.color.rgb = TEAL_MED
    av.line.width = Pt(2)
    # Initial letter
    initial = name[0]
    add_text_box(slide2, initial, cx + card_w/2 - Inches(0.65), card_y + Inches(0.3), Inches(1.3), Inches(1.3),
                 font_size=36, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
    # Name
    add_text_box(slide2, name, cx + Inches(0.15), card_y + Inches(1.75), card_w - Inches(0.3), Inches(1.1),
                 font_size=15, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    # Role pill
    pill = add_rect(slide2, cx + Inches(0.6), card_y + Inches(2.9), card_w - Inches(1.2), Inches(0.4), TEAL_LIGHT)
    add_text_box(slide2, role, cx + Inches(0.6), card_y + Inches(2.9), card_w - Inches(1.2), Inches(0.4),
                 font_size=11, color=TEAL_DARK, align=PP_ALIGN.CENTER)
 
# ─────────────────────────────────────────
# SLIDE 3 — Introdução
# ─────────────────────────────────────────
slide3 = add_slide_base(prs, LIGHT_GRAY)
add_header_bar(slide3, "Introdução")
add_footer(slide3)
 
# Left column - description
left_x = Inches(0.5)
add_text_box(slide3, "O que é o Sysplus?", left_x, Inches(1.25), Inches(5.8), Inches(0.5),
             font_size=20, bold=True, color=TEAL_DARK)
 
intro_text = (
    "O Sysplus é um marketplace online desenvolvido para conectar "
    "compradores e vendedores de forma simples, intuitiva e confiável.\n\n"
    "A plataforma oferece uma experiência completa de e-commerce, "
    "permitindo que qualquer pessoa compre ou venda produtos em múltiplas "
    "categorias — tudo em um só lugar.\n\n"
    "Com uma interface moderna e recursos personalizados, o Sysplus "
    "torna o comércio eletrônico acessível a todos."
)
add_text_box(slide3, intro_text, left_x, Inches(1.85), Inches(5.8), Inches(4.5),
             font_size=13, color=MID_GRAY, align=PP_ALIGN.LEFT)
 
# Right column - highlights box
right_x = Inches(6.8)
add_rect(slide3, right_x, Inches(1.2), Inches(6.0), Inches(5.7), WHITE)
add_rect(slide3, right_x, Inches(1.2), Inches(6.0), Inches(0.1), ORANGE_ACC)
 
add_text_box(slide3, "DESTAQUES DO PROJETO", right_x + Inches(0.3), Inches(1.4), Inches(5.4), Inches(0.4),
             font_size=11, bold=True, color=ORANGE_ACC)
 
highlights = [
    ("🛒", "Marketplace duplo", "Tanto compradores quanto vendedores têm fluxos dedicados"),
    ("🔍", "Busca inteligente", "Busca por produtos, marcas e categorias"),
    ("💳", "Múltiplos pagamentos", "Pix, cartão de crédito e débito"),
    ("🎯", "Perfil personalizado", "Recomendações por preferências e localização"),
]
for j, (icon, title, desc) in enumerate(highlights):
    by = Inches(1.95) + j * Inches(1.2)
    add_text_box(slide3, icon, right_x + Inches(0.2), by, Inches(0.5), Inches(0.5), font_size=20)
    add_text_box(slide3, title, right_x + Inches(0.8), by, Inches(4.8), Inches(0.35),
                 font_size=13, bold=True, color=DARK_GRAY)
    add_text_box(slide3, desc, right_x + Inches(0.8), by + Inches(0.35), Inches(4.8), Inches(0.5),
                 font_size=11, color=MID_GRAY)
    if j < 3:
        add_rect(slide3, right_x + Inches(0.2), by + Inches(0.95), Inches(5.6), Inches(0.02), TEAL_LIGHT)
 
# ─────────────────────────────────────────
# SLIDE 4 — Objetivos
# ─────────────────────────────────────────
slide4 = add_slide_base(prs, LIGHT_GRAY)
add_header_bar(slide4, "Objetivos")
add_footer(slide4)
 
add_text_box(slide4, "O que o Sysplus busca alcançar?", Inches(0.5), Inches(1.2), Inches(12), Inches(0.45),
             font_size=18, bold=True, color=TEAL_DARK)
 
objectives = [
    ("01", "Plataforma Acessível", "Criar um marketplace fácil de usar para compradores e vendedores de todos os perfis."),
    ("02", "Experiência Personalizada", "Oferecer recomendações e preferências baseadas no perfil e localização do usuário."),
    ("03", "Segurança nas Transações", "Garantir processos de compra e pagamento seguros e confiáveis."),
    ("04", "Capacitar Vendedores", "Fornecer ferramentas completas para que qualquer pessoa possa vender online."),
    ("05", "Interface Moderna", "Desenvolver uma UI intuitiva e responsiva com HTML, CSS e JavaScript."),
    ("06", "Gestão de Estoque", "Permitir controle de variações de produto, preço e estoque em tempo real."),
]
 
cols = 3
rows = 2
card_w2 = Inches(3.8)
card_h2 = Inches(2.1)
gap2 = Inches(0.3)
sx = (W - (cols * card_w2 + (cols-1) * gap2)) / 2
sy = Inches(1.75)
 
for idx, (num, title, desc) in enumerate(objectives):
    col = idx % cols
    row = idx // cols
    cx = sx + col * (card_w2 + gap2)
    cy = sy + row * (card_h2 + Inches(0.2))
    add_rect(slide4, cx, cy, card_w2, card_h2, WHITE)
    # num circle
    nc = slide4.shapes.add_shape(9, cx + Inches(0.2), cy + Inches(0.25), Inches(0.5), Inches(0.5))
    nc.fill.solid()
    nc.fill.fore_color.rgb = TEAL_DARK
    nc.line.fill.background()
    add_text_box(slide4, num, cx + Inches(0.2), cy + Inches(0.25), Inches(0.5), Inches(0.5),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide4, title, cx + Inches(0.8), cy + Inches(0.25), card_w2 - Inches(1.0), Inches(0.45),
                 font_size=13, bold=True, color=DARK_GRAY)
    add_text_box(slide4, desc, cx + Inches(0.2), cy + Inches(0.82), card_w2 - Inches(0.4), Inches(1.1),
                 font_size=11, color=MID_GRAY)
 
# ─────────────────────────────────────────
# SLIDE 5 — Funcionalidades
# ─────────────────────────────────────────
slide5 = add_slide_base(prs, LIGHT_GRAY)
add_header_bar(slide5, "Funcionalidades")
add_footer(slide5)
 
# Two columns: Comprador | Vendedor
col_w = Inches(5.9)
col_y = Inches(1.2)
col_h = Inches(5.7)
 
# Comprador column
add_rect(slide5, Inches(0.4), col_y, col_w, col_h, WHITE)
add_rect(slide5, Inches(0.4), col_y, col_w, Inches(0.55), TEAL_DARK)
add_text_box(slide5, "🛍️  COMPRADOR", Inches(0.4), col_y, col_w, Inches(0.55),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
 
buyer_features = [
    ("Cadastro e Login", "Conta própria com e-mail e senha"),
    ("Perfil Personalizado", "Preferências de categorias e localização por CEP"),
    ("Busca de Produtos", "Busca por produto, marca ou categoria"),
    ("Carrinho de Compras", "Adicionar, remover e visualizar total com frete"),
    ("Formas de Pagamento", "Pix (QR Code), cartão de crédito e débito"),
    ("Histórico de Pedidos", "Acompanhamento de compras realizadas"),
]
for k, (ftitle, fdesc) in enumerate(buyer_features):
    fy = col_y + Inches(0.7) + k * Inches(0.78)
    dot = slide5.shapes.add_shape(9, Inches(0.6), fy + Inches(0.12), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = TEAL_MED; dot.line.fill.background()
    add_text_box(slide5, ftitle, Inches(0.9), fy, Inches(5.0), Inches(0.3),
                 font_size=12, bold=True, color=DARK_GRAY)
    add_text_box(slide5, fdesc, Inches(0.9), fy + Inches(0.3), Inches(5.0), Inches(0.3),
                 font_size=10, color=MID_GRAY)
 
# Vendedor column
vx = Inches(7.0)
add_rect(slide5, vx, col_y, col_w, col_h, WHITE)
add_rect(slide5, vx, col_y, col_w, Inches(0.55), ORANGE_ACC)
add_text_box(slide5, "🏪  VENDEDOR", vx, col_y, col_w, Inches(0.55),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
 
seller_features = [
    ("Cadastro como Vendedor", "Criação de loja dentro da plataforma"),
    ("Cadastro de Produtos", "Título, descrição, imagens e categorias"),
    ("Variações e Estoque", "Tamanho, cor, preço e quantidade por variação"),
    ("Painel do Vendedor", "Visualização e gestão de produtos publicados"),
    ("Gerenciar Pedidos", "Acompanhamento de vendas e pedidos recebidos"),
    ("Editar e Remover", "Atualização de produtos já publicados"),
]
for k, (ftitle, fdesc) in enumerate(seller_features):
    fy = col_y + Inches(0.7) + k * Inches(0.78)
    dot = slide5.shapes.add_shape(9, vx + Inches(0.2), fy + Inches(0.12), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE_ACC; dot.line.fill.background()
    add_text_box(slide5, ftitle, vx + Inches(0.5), fy, Inches(5.0), Inches(0.3),
                 font_size=12, bold=True, color=DARK_GRAY)
    add_text_box(slide5, fdesc, vx + Inches(0.5), fy + Inches(0.3), Inches(5.0), Inches(0.3),
                 font_size=10, color=MID_GRAY)
 
# ─────────────────────────────────────────
# SLIDE 6 — Tecnologias
# ─────────────────────────────────────────
slide6 = add_slide_base(prs, LIGHT_GRAY)
add_header_bar(slide6, "Tecnologias Utilizadas")
add_footer(slide6)
 
add_text_box(slide6, "Stack de desenvolvimento do Sysplus", Inches(0.5), Inches(1.2), Inches(12), Inches(0.4),
             font_size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)
 
techs = [
    ("HTML5", "#E34F26", "Estrutura semântica de todas as páginas e componentes da interface."),
    ("CSS3", "#264DE4", "Estilização visual, responsividade e animações da interface."),
    ("JavaScript", "#F7DF1E", "Interatividade, validações, lógica de negócio no front-end e manipulação do DOM."),
    ("Figma", "#F24E1E", "Prototipação visual e design das telas antes do desenvolvimento."),
]
 
tech_w = Inches(2.8)
tech_h = Inches(3.5)
gap_t = Inches(0.38)
total_tw = len(techs) * tech_w + (len(techs) - 1) * gap_t
stx = (W - total_tw) / 2
sty = Inches(1.75)
 
hex_colors = {
    "#E34F26": RGBColor(0xE3, 0x4F, 0x26),
    "#264DE4": RGBColor(0x26, 0x4D, 0xE4),
    "#F7DF1E": RGBColor(0xF7, 0xDF, 0x1E),
    "#F24E1E": RGBColor(0xF2, 0x4E, 0x1E),
}
text_colors = {
    "#F7DF1E": DARK_GRAY,  # yellow needs dark text
}
 
for i, (name, color_hex, desc) in enumerate(techs):
    tx = stx + i * (tech_w + gap_t)
    color = hex_colors[color_hex]
    txt_color = text_colors.get(color_hex, WHITE)
 
    add_rect(slide6, tx, sty, tech_w, tech_h, WHITE)
    # Color top
    add_rect(slide6, tx, sty, tech_w, Inches(1.5), color)
    # Tech name in colored box
    add_text_box(slide6, name, tx, sty + Inches(0.5), tech_w, Inches(0.7),
                 font_size=24, bold=True, color=txt_color, align=PP_ALIGN.CENTER)
    # Description below
    add_text_box(slide6, desc, tx + Inches(0.15), sty + Inches(1.6), tech_w - Inches(0.3), Inches(1.7),
                 font_size=11, color=MID_GRAY, align=PP_ALIGN.LEFT)
 
# Bottom stat bar
add_rect(slide6, Inches(0.5), Inches(5.65), W - Inches(1.0), Inches(0.8), TEAL_DARK)
stats = [
    ("100%", "Front-End"),
    ("Semântico", "HTML estruturado"),
    ("Responsivo", "Layout adaptável"),
    ("Interativo", "DOM dinâmico"),
]
sw = (W - Inches(1.0)) / len(stats)
for i, (val, label) in enumerate(stats):
    sx2 = Inches(0.5) + i * sw
    add_text_box(slide6, val, sx2, Inches(5.65), sw, Inches(0.42),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide6, label, sx2, Inches(6.05), sw, Inches(0.3),
                 font_size=9, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)
 
# ─────────────────────────────────────────
# SLIDE 7 — Obrigado
# ─────────────────────────────────────────
slide7 = add_slide_base(prs, TEAL_DARK)
 
# Background decoration
circ2 = slide7.shapes.add_shape(9, Inches(8), Inches(-1), Inches(7), Inches(7))
circ2.fill.solid()
circ2.fill.fore_color.rgb = RGBColor(0x2E, 0x6A, 0x78)
circ2.line.fill.background()
 
circ3 = slide7.shapes.add_shape(9, Inches(-1), Inches(4), Inches(5), Inches(5))
circ3.fill.solid()
circ3.fill.fore_color.rgb = RGBColor(0x2E, 0x6A, 0x78)
circ3.line.fill.background()
 
# Orange accent bar
add_rect(slide7, Inches(0.8), Inches(2.2), Inches(0.12), Inches(2.5), ORANGE_ACC)
 
add_text_box(slide7, "Obrigado pela", Inches(1.2), Inches(2.2), Inches(8), Inches(1.0),
             font_size=44, bold=False, color=TEAL_LIGHT, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text_box(slide7, "Atenção!", Inches(1.2), Inches(3.1), Inches(8), Inches(1.1),
             font_size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
 
add_text_box(slide7, "Sysplus — Marketplace Completo", Inches(1.2), Inches(4.35), Inches(8), Inches(0.45),
             font_size=15, color=TEAL_LIGHT, align=PP_ALIGN.LEFT)
 
# Team names at bottom
team = "Artur Garcia  ·  Lucas Braga  ·  Guilherme Faria"
add_text_box(slide7, team, Inches(1.2), Inches(5.1), Inches(10), Inches(0.4),
             font_size=12, color=RGBColor(0x80, 0xBB, 0xC8), align=PP_ALIGN.LEFT)
 
# S+ logo box
logo_box = add_rect(slide7, Inches(10.2), Inches(2.8), Inches(2.2), Inches(2.2), ORANGE_ACC)
add_text_box(slide7, "S+", Inches(10.2), Inches(2.8), Inches(2.2), Inches(2.2),
             font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
 
# Save
out_path = "/home/claude/sysplus/Sysplus_Apresentacao.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
 